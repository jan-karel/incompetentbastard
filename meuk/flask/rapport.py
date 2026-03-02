"""Rapport generatie blueprint — afgesplitst van findings.py."""

from meuk.hacksec import lezen, schrijven
from flask import Blueprint, request, send_file, jsonify, render_template
from jinja2 import Environment, FileSystemLoader
from markupsafe import escape as html_escape
from meuk.flask.models import db_bevindingen, db_bevindingen_templates, db_evidence, db_notes, db, db_finding_related, db_instellingen, db_checklist
from meuk.flask.findings import _cvss_to_severity, _latex_escape, owasptop10, _get_appdata
from meuk.flask.security import require_dashboard_access
import datetime
import html as html_mod
import json
import os
import shutil
import subprocess
import threading
import uuid as _uuid

rapport_bp = Blueprint('rapport_bp', __name__,
                       template_folder='html',
                       static_folder='static')


# Async generatie state
_async_runs = {}
_async_lock = threading.Lock()


@rapport_bp.before_request
def _check_dashboard_access():
    require_dashboard_access()


# ---------------------------------------------------------------------------
# Rapport generation — LaTeX-native pipeline
# ---------------------------------------------------------------------------

_RAPPORT_DIR = os.path.join(os.path.dirname(__file__), '..', '..', 'rapport')
_HANDBOEK_DIR = os.path.join(os.path.dirname(__file__), '..', '..', 'handboek')
_EVIDENCE_SRC = os.path.join(os.path.dirname(__file__), 'db', 'evidence')

_SEVERITY_LABELS = {
    'critical': 'Kritieke',
    'high': 'Hoge',
    'medium': 'Middelhoge',
    'low': 'Lage',
    'none': 'Informatieve',
}

_SEVERITY_LABELS_EN = {
    'critical': 'Critical',
    'high': 'High',
    'medium': 'Medium',
    'low': 'Low',
    'none': 'Informational',
}

# ---------------------------------------------------------------------------
# Vertaaldict voor rapportgeneratie (NL/EN)
# ---------------------------------------------------------------------------

_T = {
    'nl': {
        'classificatie': 'Classificatie',
        'concept': 'CONCEPT',
        'inleiding': 'Inleiding',
        'management_samenvatting': 'Managementsamenvatting',
        'methodologie': 'Methodologie',
        'samenvatting': 'Samenvatting',
        'bevindingen': 'Bevindingen',
        'aanbevelingen': 'Overzicht aanbevelingen',
        'verkenning': 'Verkenning &amp; ontdekking',
        'totaal': 'Totaal',
        'aantal': 'Aantal',
        'locatie': 'Locatie',
        'impact': 'Impact',
        'aanbeveling': 'Aanbeveling',
        'legs_up': 'Legs Up',
        'startpunten': 'Assumed Breach Startpunten',
        'privilege_niveaus': 'Privilege Niveaus',
        'verstrekte_info': 'Verstrekte Informatie',
        'scenarios': "Scenario's",
        'bevindingen_label': 'bevindingen',
        'pentest_intro': 'Dit rapport bevat de bevindingen van de penetratietest uitgevoerd op {project}. Het doel was het identificeren van kwetsbaarheden die de vertrouwelijkheid, integriteit en beschikbaarheid van de systemen kunnen aantasten.',
        'redteam_intro': 'Dit rapport bevat de bevindingen van het red team assessment uitgevoerd op {project}. Het doel was het simuleren van een realistische aanval om de detectie- en responscapaciteiten van de organisatie te evalueren, evenals het identificeren van kwetsbaarheden in mensen, processen en technologie.',
        'pentest_methode': 'De penetratietest is uitgevoerd conform de OWASP Testing Guide en PTES (Penetration Testing Execution Standard). De test bestond uit de volgende fasen:',
        'redteam_methode': 'Het red team assessment is uitgevoerd conform het MITRE ATT&amp;CK framework en de Cyber Kill Chain. De test bestond uit de volgende fasen:',
        'pentest_fasen': [
            ('Verkenning', 'Informatie verzamelen over de doelomgeving'),
            ('Scanning', 'Diensten en kwetsbaarheden identificeren'),
            ('Exploitatie', 'Kwetsbaarheden bevestigen'),
            ('Post-exploitatie', 'Impact bepalen'),
            ('Rapportage', 'Documenteren'),
        ],
        'redteam_fasen': [
            ('Verkenning', 'OSINT en doelidentificatie'),
            ('Bewapening', 'Tooling en payloads voorbereiden'),
            ('Aflevering', 'Initi&euml;le toegang verkrijgen'),
            ('Exploitatie', 'Kwetsbaarheden benutten'),
            ('Installatie', 'Persistentie vestigen'),
            ('Commando &amp; Controle', 'C2-kanaal opzetten'),
            ('Acties op doel', 'Doelstellingen realiseren'),
        ],
        'scope_blackbox': 'De test is (deels) uitgevoerd als black box assessment. De testers hadden vooraf geen informatie over de interne werking, broncode of architectuur van de doelomgeving. Dit simuleert het perspectief van een externe aanvaller.',
        'scope_greybox': 'De test is (deels) uitgevoerd als grey box assessment. De testers hadden beperkte informatie, zoals netwerktopologie, gebruikersaccounts of API-documentatie. Dit simuleert een aanvaller met beperkte interne kennis.',
        'scope_whitebox': 'De test is (deels) uitgevoerd als white box assessment. De testers hadden volledige toegang tot broncode, architectuurdocumentatie en configuraties. Dit maakt een grondige analyse mogelijk van alle beveiligingslagen.',
    },
    'en': {
        'classificatie': 'Classification',
        'concept': 'DRAFT',
        'inleiding': 'Introduction',
        'management_samenvatting': 'Executive Summary',
        'methodologie': 'Methodology',
        'samenvatting': 'Summary of Findings',
        'bevindingen': 'Findings',
        'aanbevelingen': 'Recommendations Overview',
        'verkenning': 'Reconnaissance &amp; Discovery',
        'totaal': 'Total',
        'aantal': 'Count',
        'locatie': 'Location',
        'impact': 'Impact',
        'aanbeveling': 'Recommendation',
        'legs_up': 'Legs Up',
        'startpunten': 'Assumed Breach Starting Points',
        'privilege_niveaus': 'Privilege Levels',
        'verstrekte_info': 'Provided Information',
        'scenarios': 'Scenarios',
        'bevindingen_label': 'findings',
        'pentest_intro': 'This report contains the findings of the penetration test performed on {project}. The objective was to identify vulnerabilities that could compromise the confidentiality, integrity and availability of the systems.',
        'redteam_intro': 'This report contains the findings of the red team assessment performed on {project}. The objective was to simulate a realistic attack to evaluate the detection and response capabilities of the organisation, as well as to identify vulnerabilities in people, processes and technology.',
        'pentest_methode': 'The penetration test was performed in accordance with the OWASP Testing Guide and PTES (Penetration Testing Execution Standard). The test consisted of the following phases:',
        'redteam_methode': 'The red team assessment was performed in accordance with the MITRE ATT&amp;CK framework and the Cyber Kill Chain. The test consisted of the following phases:',
        'pentest_fasen': [
            ('Reconnaissance', 'Gather information about the target environment'),
            ('Scanning', 'Identify services and vulnerabilities'),
            ('Exploitation', 'Confirm vulnerabilities'),
            ('Post-exploitation', 'Determine impact'),
            ('Reporting', 'Document findings'),
        ],
        'redteam_fasen': [
            ('Reconnaissance', 'OSINT and target identification'),
            ('Weaponisation', 'Prepare tooling and payloads'),
            ('Delivery', 'Obtain initial access'),
            ('Exploitation', 'Exploit vulnerabilities'),
            ('Installation', 'Establish persistence'),
            ('Command &amp; Control', 'Set up C2 channel'),
            ('Actions on Objectives', 'Achieve goals'),
        ],
        'scope_blackbox': 'The test was (partially) performed as a black box assessment. The testers had no prior information about the internal workings, source code or architecture of the target environment. This simulates the perspective of an external attacker.',
        'scope_greybox': 'The test was (partially) performed as a grey box assessment. The testers had limited information, such as network topology, user accounts or API documentation. This simulates an attacker with limited internal knowledge.',
        'scope_whitebox': 'The test was (partially) performed as a white box assessment. The testers had full access to source code, architecture documentation and configurations. This allows for a thorough analysis of all security layers.',
    },
}

# TLP 2.0 kleuren (LaTeX kleur-naam, HTML hex, tekst-kleur voor contrast)
_TLP_MAP = {
    'TLP:RED':          ('tlpred',   '#FF2B2B', '#ffffff'),
    'TLP:AMBER':        ('tlpamber', '#FFC000', '#000000'),
    'TLP:AMBER+STRICT': ('tlpamber', '#FFC000', '#000000'),
    'TLP:GREEN':        ('tlpgreen', '#33FF00', '#000000'),
    'TLP:CLEAR':        ('tlpclear', '#000000', '#ffffff'),
}


def _copy_evidence_for_rapport(finding_id, evidence_list, dest_dir):
    """Kopieer evidence images naar rapport/evidence/ met leesbare namen."""
    copied = []
    for ev in evidence_list:
        src = os.path.join(_EVIDENCE_SRC, str(finding_id), ev.filename)
        if not os.path.exists(src):
            continue
        ct = ev.content_type or ''
        if not ct.startswith('image/'):
            continue
        safe_name = 'ev_{}_{}'.format(finding_id, ev.filename)
        dst = os.path.join(dest_dir, safe_name)
        shutil.copy2(src, dst)
        copied.append({
            'path': 'evidence/' + safe_name,
            'caption': ev.original_filename or ev.filename,
        })
    return copied


def _build_findings_data(base_query, taal='nl'):
    """Bouw een gesorteerde lijst van finding-dicts voor het LaTeX template."""
    distinct_refs = db.session.query(db_bevindingen.ref).filter(
        db_bevindingen.id.in_([b.id for b in base_query.all()])
    ).distinct().all()
    refs = [r[0] for r in distinct_refs]

    findings = []
    for ref in refs:
        template = db_bevindingen_templates.query.filter_by(id=ref).first()
        items = base_query.filter_by(ref=ref).all()
        if not items:
            continue

        for item in items:
            evidence_db = db_evidence.query.filter_by(finding_id=item.id).all()
            basescore = item.basescore or (template.basescore if template else '') or '0'
            severity, sort_key = _cvss_to_severity(basescore)

            if taal == 'en':
                beschrijving = (template.enbeschrijving if template else '') or (template.nlbeschrijving if template else '') or ''
                impact = (template.enimpact if template else '') or (template.nlimpact if template else '') or ''
                aanbeveling = (template.enaanbeveling if template else '') or (template.nlaanbeveling if template else '') or ''
            else:
                beschrijving = (template.nlbeschrijving if template else '') or ''
                impact = (template.nlimpact if template else '') or ''
                aanbeveling = (template.nlaanbeveling if template else '') or ''

            # Relaties ophalen
            rels_from = db_finding_related.query.filter_by(finding_id=item.id).all()
            related_items = []
            for r in rels_from:
                rel_bev = db.session.get(db_bevindingen, r.related_id)
                if rel_bev:
                    related_items.append({
                        'id': rel_bev.id,
                        'naam': rel_bev.naam or '',
                        'type': r.relation_type,
                    })

            findings.append({
                'id': item.id,
                'titel': item.naam or '',
                'beschrijving': beschrijving,
                'impact': impact,
                'aanbeveling': aanbeveling,
                'basescore': basescore,
                'cvss_vector': item.cvss or (template.cvss if template else '') or '',
                'locatie': item.locatie or '',
                'owasp': owasptop10[int(template.owasp) - 1] if template and template.owasp else '',
                'cwe': template.cwe if template else '',
                'severity': severity,
                'sort_key': sort_key,
                'evidence_db': evidence_db,
                'evidence': [],
                'uitwerken': item.uitwerken or '',
                'remediation_status': getattr(item, 'remediation_status', 'open') or 'open',
                'related': related_items,
            })

    # Sorteer op severity (critical=1 eerst)
    findings.sort(key=lambda f: (f['sort_key'], f['id']))
    return findings


def _build_checklist_appendix(checklist_ids, taal='nl'):
    """Bouw checklist-data voor rapport appendix.

    Returns lijst van dicts met naam, type, target, phases en stats.
    """
    from meuk.flask.checklists import _load_templates

    checklists_data = []
    templates_json = _load_templates()

    for cl in db_checklist.query.filter(db_checklist.id.in_(checklist_ids)).all():
        template = templates_json.get('checklists', {}).get(cl.checklist_type, {})
        item_map = {i.item_ref: i for i in cl.items}

        stats = {'pass': 0, 'fail': 0, 'warn': 0, 'skip': 0, 'na': 0, 'open': 0}
        phases = []
        for phase in template.get('phases', []):
            phase_items = []
            for tpl_item in phase.get('items', []):
                db_item = item_map.get(tpl_item['id'])
                status = db_item.status if db_item else 'open'
                stats[status] = stats.get(status, 0) + 1
                title = tpl_item.get('en_title' if taal == 'en' else 'title', tpl_item.get('title', ''))
                phase_items.append({
                    'ref': tpl_item['id'],
                    'title': title,
                    'status': status,
                })
            phases.append({
                'title': phase.get('title', ''),
                'items': phase_items,
            })

        total = sum(stats.values())
        en_title = template.get('en_title', template.get('title', cl.checklist_type))
        cl_title = en_title if taal == 'en' else template.get('title', cl.checklist_type)

        checklists_data.append({
            'naam': cl.naam,
            'type': cl_title,
            'target': cl.target or '',
            'phases': phases,
            'stats': stats,
            'total': total,
        })

    return checklists_data


@rapport_bp.route('/dashboard/report', methods=['GET'])
def rapport_dashboard():
    findings = db_bevindingen.query.all()
    notes = db_notes.query.order_by(db_notes.volgorde.asc(), db_notes.id.desc()).all()
    templates = db_bevindingen_templates.query.all()
    template_map = {str(t.id): t for t in templates}
    severity_map = {}
    for f in findings:
        sev, _ = _cvss_to_severity(f.basescore or (template_map.get(f.ref, None) and template_map[f.ref].basescore) or '0')
        severity_map[f.id] = sev
    checklists = db_checklist.query.order_by(db_checklist.id.desc()).all()
    return render_template('rapport_dashboard.html',
                           findings=findings, notes=notes,
                           template_map=template_map, severity_map=severity_map,
                           checklists=checklists)


@rapport_bp.route('/dashboard/report/generate', methods=['GET'])
def gen_rapport():
    os.makedirs(_RAPPORT_DIR, exist_ok=True)
    evidence_dir = os.path.join(_RAPPORT_DIR, 'evidence')
    os.makedirs(evidence_dir, exist_ok=True)

    # Status filter
    include_draft = request.args.get('include_draft', '0') == '1'
    base_query = db_bevindingen.query
    if not include_draft:
        base_query = base_query.filter(db_bevindingen.status == 'final')

    # Notities met rapport=True, gesorteerd op volgorde
    notities = db_notes.query.filter_by(rapport=True).order_by(db_notes.volgorde.asc()).all()
    notes_data = [{'naam': n.naam or '', 'uitwerken': n.uitwerken or ''} for n in notities]

    # Metadata
    s = _get_appdata()

    # Taal ophalen
    taal = (getattr(s, 'rapport_taal', '') if s else '') or 'nl'

    # Bouw findings data
    findings = _build_findings_data(base_query, taal=taal)

    # Checklists appendix
    _cl_raw = request.args.get('include_checklists', '')
    checklist_ids = [int(x) for x in _cl_raw.split(',') if x.strip().isdigit()]
    checklists_data = _build_checklist_appendix(checklist_ids, taal=taal) if checklist_ids else []

    # Kopieer evidence
    for f in findings:
        f['evidence'] = _copy_evidence_for_rapport(f['id'], f['evidence_db'], evidence_dir)

    # Groepeer per severity
    severity_order = ['critical', 'high', 'medium', 'low', 'none']
    sev_labels = _SEVERITY_LABELS_EN if taal == 'en' else _SEVERITY_LABELS
    grouped = []
    counts = {sv: 0 for sv in severity_order}
    for sev in severity_order:
        group = [f for f in findings if f['severity'] == sev]
        counts[sev] = len(group)
        grouped.append((sev_labels.get(sev, sev), group))

    # Scoped hosts
    scoped_hosts = list(set(f['locatie'] for f in findings if f['locatie']))
    titel = (s.rapport_titel if s else '') or 'Penetration Test Report'
    auteur = (s.rapport_auteur if s else '') or 'Incompetent Bastard'
    subtitel = (s.rapport_subtitel if s else '') or 'Rapportage'
    project = (getattr(s, 'rapport_project', '') if s else '') or ''
    classificatie = (getattr(s, 'rapport_classificatie', '') if s else '') or 'TLP:RED'
    is_draft = getattr(s, 'rapport_isdraft', True) if s else True
    testtype = (getattr(s, 'rapport_testtype', '') if s else '') or 'pentest'
    testscope = (getattr(s, 'rapport_testscope', '') if s else '') or 'blackbox'
    testscopes = [t.strip() for t in testscope.split(',') if t.strip()]
    management_samenvatting = (getattr(s, 'rapport_management_samenvatting', '') if s else '') or ''
    legsup_startpunt = (getattr(s, 'rapport_legsup_startpunt', '') if s else '') or ''
    legsup_privileges = (getattr(s, 'rapport_legsup_privileges', '') if s else '') or ''
    legsup_informatie = (getattr(s, 'rapport_legsup_informatie', '') if s else '') or ''
    try:
        scenarios = json.loads((getattr(s, 'rapport_scenarios', '[]') if s else '[]') or '[]')
    except (json.JSONDecodeError, TypeError):
        scenarios = []
    try:
        scope_targets = json.loads((getattr(s, 'rapport_scope_targets', '[]') if s else '[]') or '[]')
    except (json.JSONDecodeError, TypeError):
        scope_targets = []
    rapport_roe = (getattr(s, 'rapport_roe', '') if s else '') or ''
    test_start = getattr(s, 'rapport_test_start_date', None) if s else None
    test_end = getattr(s, 'rapport_test_end_date', None) if s else None
    test_start_str = test_start.strftime('%d-%m-%Y') if test_start else ''
    test_end_str = test_end.strftime('%d-%m-%Y') if test_end else ''
    variant = (getattr(s, 'rapport_template_variant', '') if s else '') or 'detailed'
    rapport_versie = (getattr(s, 'rapport_versie', '') if s else '') or ''
    rapport_versie_status = (getattr(s, 'rapport_versie_status', '') if s else '') or 'concept'
    datum = datetime.date.today().strftime('%d-%m-%Y')
    tlp_latex, tlp_html_bg, tlp_html_fg = _TLP_MAP.get(classificatie, ('tlpred', '#FF2B2B', '#ffffff'))

    # Gemeenschappelijke extra kwargs
    _extra_kwargs = dict(
        rapport_roe=rapport_roe,
        test_start=test_start_str,
        test_end=test_end_str,
        variant=variant,
        rapport_versie=rapport_versie,
        rapport_versie_status=rapport_versie_status,
    )

    # --- Render LaTeX via Jinja2 met aangepaste delimiters ---
    template_dir = os.path.join(os.path.dirname(__file__), 'html')
    latex_env = Environment(
        loader=FileSystemLoader(template_dir),
        block_start_string='[%',
        block_end_string='%]',
        variable_start_string='[[',
        variable_end_string=']]',
        comment_start_string='[#',
        comment_end_string='#]',
        autoescape=False,
    )
    latex_env.filters['latex_escape'] = _latex_escape
    tmpl = latex_env.get_template('rapport_latex.tex')

    tex_content = tmpl.render(
        titel=_latex_escape(titel),
        subtitel=_latex_escape(subtitel),
        auteur=_latex_escape(auteur),
        project=_latex_escape(project),
        classificatie=_latex_escape(classificatie),
        tlp_color=tlp_latex,
        datum=datum,
        is_draft=is_draft,
        findings=findings,
        grouped_findings=grouped,
        counts=counts,
        scoped_hosts=scoped_hosts,
        scope_targets=scope_targets,
        notes=notes_data,
        testtype=testtype,
        testscope=testscope,
        testscopes=testscopes,
        management_samenvatting=_latex_escape(management_samenvatting),
        legsup_startpunt=legsup_startpunt,
        legsup_privileges=legsup_privileges,
        legsup_informatie=legsup_informatie,
        scenarios=scenarios,
        taal=taal,
        checklists=checklists_data,
    )

    # Schrijf .tex bestand
    tex_path = os.path.join(_RAPPORT_DIR, 'rapport.tex')
    schrijven(tex_path, tex_content)

    # Legacy: schrijf ook findings_nl.tex voor backward-compat
    schrijven(os.path.join(_RAPPORT_DIR, 'findings_nl.tex'), tex_content)

    # --- PDF compilatie (optioneel) ---
    handboek_abs = os.path.abspath(_HANDBOEK_DIR)
    try:
        env = os.environ.copy()
        env['TEXINPUTS'] = '.:{handboek}:{handboek}//:{rapport}:'.format(
            handboek=handboek_abs, rapport=os.path.abspath(_RAPPORT_DIR))
        for _ in range(3):
            subprocess.run(
                ['xelatex', '-interaction=nonstopmode', '-output-directory',
                 os.path.abspath(_RAPPORT_DIR), tex_path],
                env=env, cwd=os.path.abspath(_RAPPORT_DIR),
                capture_output=True, timeout=120,
            )
    except (FileNotFoundError, subprocess.TimeoutExpired, OSError):
        pass  # xelatex niet beschikbaar

    # --- HTML generatie voor preview ---
    html = _generate_html_preview(findings, grouped, counts, titel, auteur, subtitel, datum, classificatie, is_draft, notes_data, tlp_html_bg, tlp_html_fg, management_samenvatting=management_samenvatting, testtype=testtype, testscope=testscope, testscopes=testscopes, legsup_startpunt=legsup_startpunt, legsup_privileges=legsup_privileges, legsup_informatie=legsup_informatie, scenarios=scenarios, taal=taal, scope_targets=scope_targets, checklists=checklists_data, **_extra_kwargs)
    schrijven(os.path.join(_RAPPORT_DIR, 'tex.html'), html)

    # --- STIX 2.1 bundle generatie ---
    from meuk.flask.stix_taxii import build_stix_bundle
    stix_findings = base_query.all()
    stix_templates = db_bevindingen_templates.query.all()
    stix_template_map = {str(t.id): t for t in stix_templates}
    stix = build_stix_bundle(stix_findings, stix_template_map,
                             titel=titel, auteur=auteur)
    schrijven(os.path.join(_RAPPORT_DIR, 'rapport.stix.json'),
              json.dumps(stix, indent=2, ensure_ascii=False))

    # --- DOCX Markdown generatie (altijd, onafhankelijk van pandoc) ---
    docx_md = _generate_docx_markdown(
        findings, grouped, counts, titel, auteur, subtitel, datum,
        classificatie, is_draft, notes_data,
        management_samenvatting=management_samenvatting,
        testtype=testtype, testscope=testscope, testscopes=testscopes,
        legsup_startpunt=legsup_startpunt,
        legsup_privileges=legsup_privileges,
        legsup_informatie=legsup_informatie,
        scenarios=scenarios, scoped_hosts=scoped_hosts, taal=taal,
        scope_targets=scope_targets,
        checklists=checklists_data,
        **_extra_kwargs)
    docx_md_path = os.path.join(_RAPPORT_DIR, 'rapport_docx.md')
    schrijven(docx_md_path, docx_md)

    # --- SARIF generatie ---
    sarif = _generate_sarif(findings, titel, auteur)
    schrijven(os.path.join(_RAPPORT_DIR, 'rapport.sarif.json'),
              json.dumps(sarif, indent=2, ensure_ascii=False))

    # --- PPTX generatie ---
    try:
        _generate_pptx(findings, grouped, counts, titel, auteur, subtitel, datum,
                       os.path.join(_RAPPORT_DIR, 'rapport.pptx'))
    except Exception:
        pass

    # --- MD + DOCX generatie via pandoc ---
    try:
        from sh import pandoc
        # Legacy: tex.md van HTML preview
        html_path = os.path.join(_RAPPORT_DIR, 'tex.html')
        pandoc(html_path, '-o', os.path.join(_RAPPORT_DIR, 'tex.md'))
        md = lezen(os.path.join(_RAPPORT_DIR, 'tex.md'))
        md_header = '---\ntitle: "{}"\nsubtitle: "{}"\nauthor: {}\ndate: {}\n\n...\n\n'.format(
            html_escape(titel), html_escape(subtitel), html_escape(auteur), datum)
        schrijven(os.path.join(_RAPPORT_DIR, 'tex.md'), md_header + md)
        # DOCX en ODT van dedicated Markdown (spiegelt LaTeX structuur)
        _pandoc_args = ['--toc', '--toc-depth=2',
                        '--resource-path=' + os.path.abspath(_RAPPORT_DIR)]
        pandoc(docx_md_path, *_pandoc_args,
               '-o', os.path.join(_RAPPORT_DIR, 'rapport.docx'))
        pandoc(docx_md_path, *_pandoc_args,
               '-o', os.path.join(_RAPPORT_DIR, 'rapport.odt'))
        # EPUB
        pandoc(docx_md_path, *_pandoc_args,
               '-o', os.path.join(_RAPPORT_DIR, 'rapport.epub'))
    except Exception:
        pass  # pandoc niet beschikbaar

    # --- Version snapshot ---
    _snapshot_version()

    return html


def _snapshot_version():
    """Maak een timestamped snapshot van de gegenereerde rapport bestanden."""
    ts = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
    dest = os.path.join(_RAPPORT_DIR, 'versions', ts)
    os.makedirs(dest, exist_ok=True)
    for fname in ('rapport.tex', 'rapport_docx.md', 'rapport.stix.json',
                  'rapport.sarif.json', 'tex.html'):
        src = os.path.join(_RAPPORT_DIR, fname)
        if os.path.exists(src):
            shutil.copy2(src, dest)


_STATUS_COLORS = {
    'pass': '#22c55e', 'fail': '#ef4444', 'warn': '#f59e0b',
    'skip': '#94a3b8', 'na': '#94a3b8', 'open': '#d1d5db',
}


def _append_checklist_html(parts, checklists, taal='nl'):
    """Voeg checklist appendix toe aan HTML parts."""
    if not checklists:
        return
    heading = 'Checklist Results' if taal == 'en' else 'Checklist resultaten'
    parts.append('<h2>{}</h2>'.format(heading))
    ref_l = 'Ref'
    item_l = 'Item'
    status_l = 'Status'
    for cl in checklists:
        parts.append('<h3>{} ({}) &mdash; {}</h3>'.format(
            html_escape(cl['naam']), html_escape(cl['type']), html_escape(cl['target'])))
        st = cl['stats']
        parts.append('<p><strong>{} pass, {} fail, {} warn, {} skip, {} n/a, {} open</strong></p>'.format(
            st['pass'], st['fail'], st['warn'], st['skip'], st['na'], st['open']))
        for phase in cl['phases']:
            parts.append('<h4>{}</h4>'.format(html_escape(phase['title'])))
            parts.append('<table><tr><th>{}</th><th>{}</th><th>{}</th></tr>'.format(ref_l, item_l, status_l))
            for item in phase['items']:
                color = _STATUS_COLORS.get(item['status'], '#d1d5db')
                parts.append('<tr><td><code>{}</code></td><td>{}</td><td><span style="color:{}">{}</span></td></tr>'.format(
                    html_escape(item['ref']), html_escape(item['title']),
                    color, html_escape(item['status'].upper())))
            parts.append('</table>')


def _append_signoff_html(parts, taal='nl'):
    """Voeg sign-off tabel toe aan HTML parts."""
    label = 'Sign-off'
    role_l = 'Role' if taal == 'en' else 'Rol'
    name_l = 'Name' if taal == 'en' else 'Naam'
    date_l = 'Date' if taal == 'en' else 'Datum'
    sig_l = 'Signature' if taal == 'en' else 'Handtekening'
    parts.append('<h2>{}</h2>'.format(label))
    parts.append('<table><tr><th>{}</th><th>{}</th><th>{}</th><th>{}</th></tr>'.format(role_l, name_l, date_l, sig_l))
    for role in (['Lead Tester', 'Reviewer', 'Project Manager'] if taal == 'en' else ['Lead Tester', 'Reviewer', 'Projectmanager']):
        parts.append('<tr><td>{}</td><td>&nbsp;</td><td>&nbsp;</td><td>&nbsp;</td></tr>'.format(role))
    parts.append('</table>')


def _generate_html_preview(findings, grouped, counts, titel, auteur, subtitel, datum, classificatie, is_draft, notes=None, tlp_bg='#FF2B2B', tlp_fg='#ffffff', management_samenvatting='', testtype='pentest', testscope='blackbox', testscopes=None, legsup_startpunt='', legsup_privileges='', legsup_informatie='', scenarios=None, taal='nl', rapport_roe='', test_start='', test_end='', variant='detailed', rapport_versie='', rapport_versie_status='concept', scope_targets=None, checklists=None):
    """Genereer HTML preview van het rapport."""
    t = _T.get(taal, _T['nl'])
    sev_labels = _SEVERITY_LABELS_EN if taal == 'en' else _SEVERITY_LABELS
    severity_colors = {
        'critical': '#820000', 'high': '#e60000', 'medium': '#f09600',
        'low': '#ffdc00', 'none': '#00be00',
    }
    parts = ['<html><head><style>',
             'body{font-family:sans-serif;max-width:900px;margin:0 auto;padding:20px}',
             'pre{background:#f0f0f0;padding:10px;overflow-x:auto}',
             'img{max-width:100%}',
             'table{width:100%;border-collapse:collapse;margin:1em 0}',
             'th,td{border:1px solid #ccc;padding:6px 10px;text-align:left}',
             'th{background:#f0f0f0}',
             '.severity-bar{width:4px;display:inline-block;height:1em;margin-right:6px}',
             '.finding{border-left:4px solid #ccc;padding:12px;margin:12px 0}',
             '.tlp-badge{display:inline-block;padding:2px 10px;border-radius:4px;font-weight:bold;font-size:0.9em}',
             '</style></head><body>']

    # Titel + versie + TLP badge
    parts.append('<h1>{}</h1>'.format(html_escape(titel)))
    parts.append('<p><strong>{}</strong></p>'.format(html_escape(subtitel)))
    if rapport_versie:
        _vs_labels = {'concept': 'Concept', 'definitief': 'Definitief', 'herzien': 'Herzien'}
        _vs_label = _vs_labels.get(rapport_versie_status, rapport_versie_status.capitalize())
        parts.append('<p>Versie {} &mdash; {}</p>'.format(html_escape(rapport_versie), html_escape(_vs_label)))
    parts.append('<p>{} &mdash; {}</p>'.format(html_escape(auteur), html_escape(datum)))
    parts.append('<p>{}: <span class="tlp-badge" style="background:{};color:{}">{}</span></p>'.format(
        t['classificatie'], tlp_bg, tlp_fg, html_escape(classificatie)))
    if is_draft:
        parts.append('<p style="color:red;font-size:1.2em"><strong>{}</strong></p>'.format(t['concept']))

    # Board variant: KPI + top 3 only
    if variant == 'board':
        parts.append('<h2>Summary</h2>' if taal == 'en' else '<h2>Samenvatting</h2>')
        parts.append('<table><tr><th>{}</th><th>{}</th></tr>'.format(t['classificatie'], t['aantal']))
        for sev in ('critical', 'high', 'medium', 'low', 'none'):
            color = severity_colors[sev]
            label = sev_labels.get(sev, sev)
            parts.append('<tr><td><span class="severity-bar" style="background:{}"></span>{}</td><td>{}</td></tr>'.format(
                color, html_escape(label), counts[sev]))
        parts.append('<tr><td><strong>{}</strong></td><td><strong>{}</strong></td></tr></table>'.format(t['totaal'], len(findings)))
        # Top 3 findings
        top3 = findings[:3]
        if top3:
            parts.append('<h2>Top Findings</h2>' if taal == 'en' else '<h2>Top Bevindingen</h2>')
            for f in top3:
                color = severity_colors.get(f['severity'], '#ccc')
                parts.append('<div class="finding" style="border-left-color:{}">'.format(color))
                parts.append('<h4>[{}] {}</h4>'.format(str(f['id']).zfill(3), html_escape(f['titel'])))
                parts.append('<p><strong>CVSS:</strong> {} ({})</p>'.format(
                    html_escape(str(f['basescore'])), html_escape(f['severity'])))
                if f['aanbeveling']:
                    parts.append('<p><strong>{}:</strong> {}</p>'.format(t['aanbeveling'], html_escape(f['aanbeveling'])))
                parts.append('</div>')
        parts.append('</body></html>')
        return ''.join(parts)

    # Inleiding
    parts.append('<h2>{}</h2>'.format(t['inleiding']))
    intro_key = 'redteam_intro' if testtype == 'redteam' else 'pentest_intro'
    parts.append('<p>{}</p>'.format(t[intro_key].format(project=html_escape(titel))))
    parts.append('<p>{} &mdash; {}</p>'.format(html_escape(auteur), html_escape(datum)))

    # Test timeline
    if test_start or test_end:
        tl_label = 'Test Period' if taal == 'en' else 'Testperiode'
        parts.append('<p><strong>{}:</strong> {} — {}</p>'.format(
            tl_label, html_escape(test_start or '?'), html_escape(test_end or '?')))

    # Rules of Engagement
    if rapport_roe:
        roe_label = 'Rules of Engagement'
        parts.append('<h2>{}</h2>'.format(roe_label))
        parts.append('<p>{}</p>'.format(html_escape(rapport_roe)))

    # Scope targets
    if scope_targets:
        scope_label = 'Scope' if taal == 'en' else 'Scope'
        target_l = 'Target'
        type_l = 'Type'
        desc_l = 'Description' if taal == 'en' else 'Beschrijving'
        parts.append('<h2>{}</h2>'.format(scope_label))
        parts.append('<table><tr><th>{}</th><th>{}</th><th>{}</th></tr>'.format(target_l, type_l, desc_l))
        for st in scope_targets:
            if isinstance(st, dict):
                parts.append('<tr><td>{}</td><td>{}</td><td>{}</td></tr>'.format(
                    html_escape(st.get('target', '')),
                    html_escape(st.get('type', '')),
                    html_escape(st.get('beschrijving', ''))))
        parts.append('</table>')

    # Managementsamenvatting
    if management_samenvatting:
        parts.append('<h2>{}</h2>'.format(t['management_samenvatting']))
        parts.append('<p>{}</p>'.format(html_escape(management_samenvatting)))

    # Methodologie
    parts.append('<h2>{}</h2>'.format(t['methodologie']))
    methode_key = 'redteam_methode' if testtype == 'redteam' else 'pentest_methode'
    fasen_key = 'redteam_fasen' if testtype == 'redteam' else 'pentest_fasen'
    parts.append('<p>{}</p>'.format(t[methode_key]))
    parts.append('<ol>')
    for fase_naam, fase_beschrijving in t[fasen_key]:
        parts.append('<li><strong>{}</strong> &mdash; {}</li>'.format(fase_naam, fase_beschrijving))
    parts.append('</ol>')

    # Scope type (multi-select)
    scopes = testscopes if testscopes else [testscope or 'blackbox']
    for scope in scopes:
        if scope == 'whitebox':
            parts.append('<p>{}</p>'.format(t['scope_whitebox']))
        elif scope == 'greybox':
            parts.append('<p>{}</p>'.format(t['scope_greybox']))
        else:
            parts.append('<p>{}</p>'.format(t['scope_blackbox']))

    # Legs up (alleen bij redteam)
    if testtype == 'redteam' and (legsup_startpunt or legsup_privileges or legsup_informatie):
        parts.append('<h2>{}</h2>'.format(t['legs_up']))
        if legsup_startpunt:
            parts.append('<h3>{}</h3>'.format(t['startpunten']))
            parts.append('<p>{}</p>'.format(html_escape(legsup_startpunt)))
        if legsup_privileges:
            parts.append('<h3>{}</h3>'.format(t['privilege_niveaus']))
            parts.append('<p>{}</p>'.format(html_escape(legsup_privileges)))
        if legsup_informatie:
            parts.append('<h3>{}</h3>'.format(t['verstrekte_info']))
            parts.append('<p>{}</p>'.format(html_escape(legsup_informatie)))

    # Scenario's (alleen bij redteam)
    if testtype == 'redteam' and scenarios:
        parts.append('<h2>{}</h2>'.format(t['scenarios']))
        for sc in scenarios:
            naam = sc.get('naam', '') if isinstance(sc, dict) else ''
            beschrijving = sc.get('beschrijving', '') if isinstance(sc, dict) else ''
            if naam or beschrijving:
                parts.append('<h3>{}</h3>'.format(html_escape(naam)))
                parts.append('<p>{}</p>'.format(html_escape(beschrijving)))

    # Samenvatting
    parts.append('<h2>{}</h2>'.format(t['samenvatting']))
    parts.append('<table><tr><th>{}</th><th>{}</th></tr>'.format(t['classificatie'], t['aantal']))
    for sev in ('critical', 'high', 'medium', 'low', 'none'):
        color = severity_colors[sev]
        label = sev_labels.get(sev, sev)
        parts.append('<tr><td><span class="severity-bar" style="background:{}"></span>{}</td><td>{}</td></tr>'.format(
            color, html_escape(label), counts[sev]))
    parts.append('<tr><td><strong>{}</strong></td><td><strong>{}</strong></td></tr></table>'.format(t['totaal'], len(findings)))

    # Executive variant: KPI + critical/high + recommendations only
    if variant == 'executive':
        exec_findings = [f for f in findings if f['severity'] in ('critical', 'high')]
        if exec_findings:
            parts.append('<h2>{}</h2>'.format(t['bevindingen']))
            for f in exec_findings:
                color = severity_colors.get(f['severity'], '#ccc')
                parts.append('<div class="finding" style="border-left-color:{}">'.format(color))
                parts.append('<h4>[{}] {}</h4>'.format(str(f['id']).zfill(3), html_escape(f['titel'])))
                parts.append('<p><strong>CVSS:</strong> {} ({})</p>'.format(
                    html_escape(str(f['basescore'])), html_escape(f['severity'])))
                if f['aanbeveling']:
                    parts.append('<p><strong>{}:</strong> {}</p>'.format(t['aanbeveling'], html_escape(f['aanbeveling'])))
                parts.append('</div>')
        # Aanbevelingen
        parts.append('<h2>{}</h2>'.format(t['aanbevelingen']))
        parts.append('<table><tr><th>ID</th><th>{}</th><th>CVSS</th><th>{}</th></tr>'.format(t['bevindingen'], t['aanbeveling']))
        for f in findings:
            parts.append('<tr><td>{}</td><td>{}</td><td>{}</td><td>{}</td></tr>'.format(
                str(f['id']).zfill(3), html_escape(f['titel']),
                html_escape(str(f['basescore'])), html_escape(f['aanbeveling'])))
        parts.append('</table>')
        # Sign-off
        _append_signoff_html(parts, taal)
        parts.append('</body></html>')
        return ''.join(parts)

    # Notities (verkenning & ontdekking)
    if notes:
        parts.append('<h2>{}</h2>'.format(t['verkenning']))
        for n in notes:
            parts.append('<h3>{}</h3>'.format(html_escape(n['naam'])))
            parts.append('<p>{}</p>'.format(html_escape(n['uitwerken'])))

    # Bevindingen (detailed variant — full content)
    _REM_LABELS = {'open': 'Open', 'in_progress': 'In Progress', 'fixed': 'Fixed', 'verified': 'Verified', 'accepted': 'Accepted'}
    parts.append('<h2>{}</h2>'.format(t['bevindingen']))
    for label, group in grouped:
        if not group:
            continue
        parts.append('<h3>{} {}</h3>'.format(html_escape(label), t['bevindingen_label']))
        for f in group:
            color = severity_colors.get(f['severity'], '#ccc')
            parts.append('<div class="finding" style="border-left-color:{}">'.format(color))
            parts.append('<h4>[{}] {}</h4>'.format(str(f['id']).zfill(3), html_escape(f['titel'])))
            parts.append('<p><strong>CVSS:</strong> {} ({})</p>'.format(
                html_escape(str(f['basescore'])), html_escape(f['severity'])))
            # Remediation badge
            rem = f.get('remediation_status', 'open')
            rem_label = _REM_LABELS.get(rem, rem)
            parts.append('<p><strong>Remediation:</strong> <span style="padding:2px 8px;border-radius:3px;background:#eee">{}</span></p>'.format(html_escape(rem_label)))
            if f['locatie']:
                parts.append('<p><strong>{}:</strong> {}</p>'.format(t['locatie'], html_escape(f['locatie'])))
            if f['beschrijving']:
                parts.append('<p>{}</p>'.format(html_escape(f['beschrijving'])))
            if f['impact']:
                parts.append('<p><strong>{}:</strong> {}</p>'.format(t['impact'], html_escape(f['impact'])))
            if f['aanbeveling']:
                parts.append('<p><strong>{}:</strong> {}</p>'.format(t['aanbeveling'], html_escape(f['aanbeveling'])))
            # Finding relaties
            related = f.get('related', [])
            if related:
                rel_parts = []
                for r in related:
                    rel_parts.append('[{}] {} ({})'.format(r['type'], html_escape(r['naam']), r['id']))
                parts.append('<p><strong>Related:</strong> {}</p>'.format(', '.join(rel_parts)))
            for ev in f.get('evidence', []):
                parts.append('<p><em>Evidence: {}</em></p>'.format(html_escape(ev['caption'])))
            parts.append('</div>')

    # Aanbevelingen tabel
    parts.append('<h2>{}</h2>'.format(t['aanbevelingen']))
    parts.append('<table><tr><th>ID</th><th>{}</th><th>CVSS</th><th>{}</th></tr>'.format(t['bevindingen'], t['aanbeveling']))
    for f in findings:
        parts.append('<tr><td>{}</td><td>{}</td><td>{}</td><td>{}</td></tr>'.format(
            str(f['id']).zfill(3), html_escape(f['titel']),
            html_escape(str(f['basescore'])), html_escape(f['aanbeveling'])))
    parts.append('</table>')

    # Checklist appendix
    _append_checklist_html(parts, checklists, taal)

    # Sign-off
    _append_signoff_html(parts, taal)

    parts.append('</body></html>')
    return ''.join(parts)


def _generate_docx_markdown(findings, grouped, counts, titel, auteur, subtitel,
                            datum, classificatie, is_draft, notes=None,
                            management_samenvatting='', testtype='pentest',
                            testscope='blackbox', testscopes=None,
                            legsup_startpunt='', legsup_privileges='',
                            legsup_informatie='', scenarios=None,
                            scoped_hosts=None, taal='nl',
                            rapport_roe='', test_start='', test_end='',
                            variant='detailed', rapport_versie='',
                            rapport_versie_status='concept',
                            scope_targets=None, checklists=None):
    """Genereer pandoc-Markdown dat de LaTeX template 1:1 spiegelt."""
    t = _T.get(taal, _T['nl'])
    sev_labels = _SEVERITY_LABELS_EN if taal == 'en' else _SEVERITY_LABELS

    def u(text):
        """Unescape HTML entities voor Markdown output."""
        return html_mod.unescape(text)

    lines = []

    # YAML metadata block
    lines.append('---')
    lines.append('title: "{}"'.format(titel.replace('"', '\\"')))
    lines.append('subtitle: "{}"'.format(subtitel.replace('"', '\\"')))
    lines.append('author: "{}"'.format(auteur.replace('"', '\\"')))
    lines.append('date: "{}"'.format(datum))
    lines.append('---')
    lines.append('')

    if rapport_versie:
        _vs_labels = {'concept': 'Concept', 'definitief': 'Definitief', 'herzien': 'Herzien'}
        _vs_label = _vs_labels.get(rapport_versie_status, rapport_versie_status.capitalize())
        lines.append('**Versie {} --- {}**'.format(rapport_versie, _vs_label))
        lines.append('')

    if is_draft:
        lines.append('**{}**'.format(u(t['concept'])))
        lines.append('')

    # Inleiding
    lines.append('# {}'.format(u(t['inleiding'])))
    lines.append('')
    intro_key = 'redteam_intro' if testtype == 'redteam' else 'pentest_intro'
    lines.append(u(t[intro_key]).format(project=titel))
    lines.append('')
    lines.append('{} — {}'.format(auteur, datum))
    lines.append('')

    # Test timeline
    if test_start or test_end:
        tl_label = 'Test Period' if taal == 'en' else 'Testperiode'
        lines.append('**{}:** {} — {}'.format(tl_label, test_start or '?', test_end or '?'))
        lines.append('')

    # Rules of Engagement
    if rapport_roe:
        lines.append('## Rules of Engagement')
        lines.append('')
        lines.append(rapport_roe)
        lines.append('')

    # Scope subsectie
    lines.append('## Scope')
    lines.append('')
    if taal == 'en':
        lines.append('The following systems were part of the test:')
    else:
        lines.append('De volgende systemen zijn onderdeel van de test:')
    lines.append('')
    if scope_targets:
        target_l = 'Target'
        type_l = 'Type'
        desc_l = 'Description' if taal == 'en' else 'Beschrijving'
        lines.append('| {} | {} | {} |'.format(target_l, type_l, desc_l))
        lines.append('|---|---|---|')
        for st in scope_targets:
            if isinstance(st, dict):
                lines.append('| {} | {} | {} |'.format(
                    st.get('target', ''), st.get('type', ''), st.get('beschrijving', '')))
    elif scoped_hosts:
        for host in scoped_hosts:
            lines.append('- {}'.format(host))
    else:
        if taal == 'en':
            lines.append('- See findings for specific locations')
        else:
            lines.append('- Zie bevindingen voor specifieke locaties')
    lines.append('')

    # Managementsamenvatting
    if management_samenvatting:
        lines.append('# {}'.format(u(t['management_samenvatting'])))
        lines.append('')
        lines.append(management_samenvatting)
        lines.append('')

    # Methodologie
    lines.append('# {}'.format(u(t['methodologie'])))
    lines.append('')
    methode_key = 'redteam_methode' if testtype == 'redteam' else 'pentest_methode'
    fasen_key = 'redteam_fasen' if testtype == 'redteam' else 'pentest_fasen'
    lines.append(u(t[methode_key]))
    lines.append('')
    for i, (fase_naam, fase_beschrijving) in enumerate(t[fasen_key], 1):
        lines.append('{}. **{}** — {}'.format(i, u(fase_naam), u(fase_beschrijving)))
    lines.append('')

    # Scope type paragrafen
    scopes = testscopes if testscopes else [testscope or 'blackbox']
    for scope in scopes:
        if scope == 'whitebox':
            lines.append(u(t['scope_whitebox']))
        elif scope == 'greybox':
            lines.append(u(t['scope_greybox']))
        else:
            lines.append(u(t['scope_blackbox']))
        lines.append('')

    # Legs up (alleen redteam)
    if testtype == 'redteam' and (legsup_startpunt or legsup_privileges or legsup_informatie):
        lines.append('# {}'.format(u(t['legs_up'])))
        lines.append('')
        if legsup_startpunt:
            lines.append('## {}'.format(u(t['startpunten'])))
            lines.append('')
            lines.append(legsup_startpunt)
            lines.append('')
        if legsup_privileges:
            lines.append('## {}'.format(u(t['privilege_niveaus'])))
            lines.append('')
            lines.append(legsup_privileges)
            lines.append('')
        if legsup_informatie:
            lines.append('## {}'.format(u(t['verstrekte_info'])))
            lines.append('')
            lines.append(legsup_informatie)
            lines.append('')

    # Scenario's (alleen redteam)
    if testtype == 'redteam' and scenarios:
        lines.append('# {}'.format(u(t['scenarios'])))
        lines.append('')
        for sc in scenarios:
            naam = sc.get('naam', '') if isinstance(sc, dict) else ''
            beschrijving = sc.get('beschrijving', '') if isinstance(sc, dict) else ''
            if naam or beschrijving:
                lines.append('## {}'.format(naam))
                lines.append('')
                lines.append(beschrijving)
                lines.append('')

    # Risicoinschatting
    if taal == 'en':
        lines.append('# Risk Assessment')
        lines.append('')
        lines.append('The severity of findings is determined using the CVSS 4.0 framework '
                     '(Common Vulnerability Scoring System). The following classification is used:')
        lines.append('')
        lines.append('| Classification | CVSS Score | Colour |')
        lines.append('|---|---|---|')
        lines.append('| Critical | 9.0 -- 10.0 | Dark red |')
        lines.append('| High | 7.0 -- 8.9 | Red |')
        lines.append('| Medium | 4.0 -- 6.9 | Orange |')
        lines.append('| Low | 0.1 -- 3.9 | Yellow |')
        lines.append('| Informational | 0.0 | Green |')
    else:
        lines.append('# Risicoinschatting')
        lines.append('')
        lines.append('De ernst van bevindingen wordt bepaald aan de hand van het CVSS 4.0 framework '
                     '(Common Vulnerability Scoring System). De volgende classificatie wordt gehanteerd:')
        lines.append('')
        lines.append('| Classificatie | CVSS Score | Kleur |')
        lines.append('|---|---|---|')
        lines.append('| Kritiek | 9.0 -- 10.0 | Donkerrood |')
        lines.append('| Hoog | 7.0 -- 8.9 | Rood |')
        lines.append('| Midden | 4.0 -- 6.9 | Oranje |')
        lines.append('| Laag | 0.1 -- 3.9 | Geel |')
        lines.append('| Nihil | 0.0 | Groen |')
    lines.append('')

    # Samenvatting bevindingen
    lines.append('# {}'.format(u(t['samenvatting'])))
    lines.append('')
    if taal == 'en':
        lines.append('A total of **{}** findings were identified:'.format(len(findings)))
    else:
        lines.append('In totaal zijn **{}** bevindingen ge\u00efdentificeerd:'.format(len(findings)))
    lines.append('')
    lines.append('| {} | {} |'.format(u(t['classificatie']), u(t['aantal'])))
    lines.append('|---|---|')
    for sev in ('critical', 'high', 'medium', 'low', 'none'):
        label = sev_labels.get(sev, sev)
        lines.append('| {} | {} |'.format(label, counts[sev]))
    lines.append('| **{}** | **{}** |'.format(u(t['totaal']), len(findings)))
    lines.append('')

    # Verkenning & ontdekking (notes)
    if notes:
        lines.append('# {}'.format(u(t['verkenning'])))
        lines.append('')
        for n in notes:
            lines.append('## {}'.format(n['naam']))
            lines.append('')
            lines.append(n['uitwerken'])
            lines.append('')

    # Bevindingen
    _REM_LABELS = {'open': 'Open', 'in_progress': 'In Progress', 'fixed': 'Fixed', 'verified': 'Verified', 'accepted': 'Accepted'}
    lines.append('# {}'.format(u(t['bevindingen'])))
    lines.append('')

    # Board variant: top 3 only
    if variant == 'board':
        top3 = findings[:3]
        for f in top3:
            lines.append('### [{}] {} (CVSS: {})'.format(
                str(f['id']).zfill(3), f['titel'], f['basescore']))
            lines.append('')
            if f['aanbeveling']:
                lines.append('**{}:** {}'.format(u(t['aanbeveling']), f['aanbeveling']))
                lines.append('')
        return '\n'.join(lines)

    # Executive variant: critical/high only
    if variant == 'executive':
        exec_findings = [f for f in findings if f['severity'] in ('critical', 'high')]
        for f in exec_findings:
            lines.append('### [{}] {} (CVSS: {})'.format(
                str(f['id']).zfill(3), f['titel'], f['basescore']))
            lines.append('')
            if f['aanbeveling']:
                lines.append('**{}:** {}'.format(u(t['aanbeveling']), f['aanbeveling']))
                lines.append('')
    else:
        # Detailed variant: alle findings met volledige info
        for label, group in grouped:
            if not group:
                continue
            lines.append('## {} {}'.format(label, u(t['bevindingen_label'])))
            lines.append('')
            for f in group:
                lines.append('### [{}] {} (CVSS: {})'.format(
                    str(f['id']).zfill(3), f['titel'], f['basescore']))
                lines.append('')
                # Remediation badge
                rem = f.get('remediation_status', 'open')
                lines.append('**Remediation:** {}'.format(_REM_LABELS.get(rem, rem)))
                lines.append('')
                if f['beschrijving']:
                    lines.append(f['beschrijving'])
                    lines.append('')
                if f['impact']:
                    lines.append('**{}:** {}'.format(u(t['impact']), f['impact']))
                    lines.append('')
                if f['aanbeveling']:
                    lines.append('**{}:** {}'.format(u(t['aanbeveling']), f['aanbeveling']))
                    lines.append('')
                if f['locatie']:
                    lines.append('**{}:** `{}`'.format(u(t['locatie']), f['locatie']))
                    lines.append('')
                # Metadata: OWASP / CWE / CVSS Vector
                meta_parts = []
                if f['owasp']:
                    meta_parts.append('**OWASP:** {}'.format(f['owasp']))
                if f['cwe']:
                    meta_parts.append('**CWE:** {}'.format(f['cwe']))
                if f['cvss_vector']:
                    meta_parts.append('**CVSS Vector:** `{}`'.format(f['cvss_vector']))
                if meta_parts:
                    lines.append(' | '.join(meta_parts))
                    lines.append('')
                # Finding relaties
                related = f.get('related', [])
                if related:
                    rel_parts = []
                    for r in related:
                        rel_parts.append('[{}] {} ({})'.format(r['type'], r['naam'], r['id']))
                    lines.append('**Related:** {}'.format(', '.join(rel_parts)))
                    lines.append('')
                # Evidence afbeeldingen
                for ev in f.get('evidence', []):
                    lines.append('![{}]({})'.format(ev['caption'], ev['path']))
                    lines.append('')

    # Overzicht aanbevelingen
    lines.append('# {}'.format(u(t['aanbevelingen'])))
    lines.append('')
    lines.append('| ID | {} | CVSS | {} |'.format(u(t['bevindingen']), u(t['aanbeveling'])))
    lines.append('|---|---|---|---|')
    for f in findings:
        lines.append('| {} | {} | {} | {} |'.format(
            str(f['id']).zfill(3), f['titel'], f['basescore'], f['aanbeveling']))
    lines.append('')

    # Checklist appendix
    if checklists:
        heading = 'Checklist Results' if taal == 'en' else 'Checklist resultaten'
        lines.append('# {}'.format(heading))
        lines.append('')
        for cl in checklists:
            lines.append('## {} ({}) — {}'.format(cl['naam'], cl['type'], cl['target']))
            lines.append('')
            st = cl['stats']
            lines.append('**{} pass, {} fail, {} warn, {} skip, {} n/a, {} open**'.format(
                st['pass'], st['fail'], st['warn'], st['skip'], st['na'], st['open']))
            lines.append('')
            for phase in cl['phases']:
                lines.append('### {}'.format(phase['title']))
                lines.append('')
                lines.append('| Ref | Item | Status |')
                lines.append('|---|---|---|')
                for item in phase['items']:
                    lines.append('| {} | {} | {} |'.format(
                        item['ref'], item['title'], item['status'].upper()))
                lines.append('')

    # Sign-off sectie
    lines.append('# Sign-off')
    lines.append('')
    role_l = 'Role' if taal == 'en' else 'Rol'
    name_l = 'Name' if taal == 'en' else 'Naam'
    date_l = 'Date' if taal == 'en' else 'Datum'
    sig_l = 'Signature' if taal == 'en' else 'Handtekening'
    lines.append('| {} | {} | {} | {} |'.format(role_l, name_l, date_l, sig_l))
    lines.append('|---|---|---|---|')
    for role in (['Lead Tester', 'Reviewer', 'Project Manager'] if taal == 'en' else ['Lead Tester', 'Reviewer', 'Projectmanager']):
        lines.append('| {} | | | |'.format(role))
    lines.append('')

    return '\n'.join(lines)


# ---------------------------------------------------------------------------
# SARIF 2.1.0 generatie
# ---------------------------------------------------------------------------

def _generate_sarif(findings, titel='', auteur=''):
    """Genereer SARIF 2.1.0 JSON output."""
    rules = []
    results = []
    for f in findings:
        rule_id = 'FINDING-{}'.format(str(f['id']).zfill(3))
        rules.append({
            'id': rule_id,
            'name': f['titel'],
            'shortDescription': {'text': f['titel']},
            'fullDescription': {'text': f.get('beschrijving', '') or f['titel']},
            'properties': {
                'cvss': f['basescore'],
                'severity': f['severity'],
            },
        })
        result = {
            'ruleId': rule_id,
            'level': 'error' if f['severity'] in ('critical', 'high') else 'warning' if f['severity'] == 'medium' else 'note',
            'message': {'text': f.get('beschrijving', '') or f['titel']},
            'properties': {
                'cvss_score': f['basescore'],
                'cvss_vector': f.get('cvss_vector', ''),
                'remediation_status': f.get('remediation_status', 'open'),
            },
        }
        if f.get('locatie'):
            result['locations'] = [{'physicalLocation': {'artifactLocation': {'uri': f['locatie']}}}]
        if f.get('aanbeveling'):
            result['fixes'] = [{'description': {'text': f['aanbeveling']}}]
        results.append(result)

    return {
        '$schema': 'https://raw.githubusercontent.com/oasis-tcs/sarif-spec/master/Schemata/sarif-schema-2.1.0.json',
        'version': '2.1.0',
        'runs': [{
            'tool': {
                'driver': {
                    'name': 'Incompetent Bastard',
                    'version': '1.0',
                    'informationUri': 'https://github.com/incompetentbastard',
                    'rules': rules,
                },
            },
            'results': results,
            'invocations': [{
                'executionSuccessful': True,
                'properties': {
                    'report_title': titel,
                    'author': auteur,
                },
            }],
        }],
    }


# ---------------------------------------------------------------------------
# PPTX generatie
# ---------------------------------------------------------------------------

def _generate_pptx(findings, grouped, counts, titel, auteur, subtitel, datum, output_path):
    """Genereer PowerPoint presentatie met KPI's en findings."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Titelpagina
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = titel
    if slide.placeholders[1]:
        slide.placeholders[1].text = '{}\n{}\n{}'.format(subtitel, auteur, datum)

    # Slide 2: KPI overview
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = 'Summary'
    left = Inches(1)
    top = Inches(2)
    width = Inches(11)
    height = Inches(4)
    table = slide.shapes.add_table(7, 2, left, top, width, height).table
    table.cell(0, 0).text = 'Severity'
    table.cell(0, 1).text = 'Count'
    for i, sev in enumerate(['critical', 'high', 'medium', 'low', 'none'], 1):
        table.cell(i, 0).text = sev.capitalize()
        table.cell(i, 1).text = str(counts.get(sev, 0))
    table.cell(6, 0).text = 'Total'
    table.cell(6, 1).text = str(len(findings))

    # Slides per finding (top 10 max)
    for f in findings[:10]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = '[{}] {}'.format(str(f['id']).zfill(3), f['titel'])
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = 'CVSS: {} ({})'.format(f['basescore'], f['severity'])
        if f.get('beschrijving'):
            p = tf.add_paragraph()
            p.text = f['beschrijving'][:500]
            p.font.size = Pt(14)
        if f.get('aanbeveling'):
            p = tf.add_paragraph()
            p.text = 'Recommendation: ' + f['aanbeveling'][:300]
            p.font.size = Pt(14)
            p.font.bold = True

    prs.save(output_path)


# ---------------------------------------------------------------------------
# Rapport download
# ---------------------------------------------------------------------------

@rapport_bp.route('/api/report/download/<fmt>')
def rapport_download(fmt):
    FORMATS = {
        'html': ('tex.html', 'text/html'),
        'md':   ('tex.md',   'text/markdown'),
        'tex':  ('rapport.tex', 'application/x-tex'),
        'pdf':  ('rapport.pdf', 'application/pdf'),
        'docx': ('rapport.docx', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'),
        'odt':  ('rapport.odt',  'application/vnd.oasis.opendocument.text'),
        'stix': ('rapport.stix.json', 'application/json'),
        'sarif': ('rapport.sarif.json', 'application/json'),
        'pptx': ('rapport.pptx', 'application/vnd.openxmlformats-officedocument.presentationml.presentation'),
        'epub': ('rapport.epub', 'application/epub+zip'),
    }
    if fmt not in FORMATS:
        return jsonify({'error': 'Ongeldig formaat'}), 400
    filename, mimetype = FORMATS[fmt]
    path = os.path.join(_RAPPORT_DIR, filename)
    if not os.path.exists(path):
        return jsonify({'error': 'Genereer eerst het rapport'}), 404
    return send_file(path, mimetype=mimetype, as_attachment=True,
                     download_name=f'rapport.{fmt}')


# ---------------------------------------------------------------------------
# Logo upload: POST /api/report/logo
# ---------------------------------------------------------------------------

@rapport_bp.route('/api/report/logo', methods=['POST'])
def rapport_logo_upload():
    from werkzeug.utils import secure_filename as _sf
    if 'file' not in request.files:
        return jsonify({'ok': False, 'error': 'geen bestand'}), 400
    f = request.files['file']
    if not f.filename:
        return jsonify({'ok': False, 'error': 'geen bestandsnaam'}), 400
    ct = f.content_type or ''
    if not ct.startswith('image/'):
        return jsonify({'ok': False, 'error': 'alleen afbeeldingen toegestaan'}), 400
    os.makedirs(_RAPPORT_DIR, exist_ok=True)
    safe = _sf(f.filename)
    dest = os.path.join(_RAPPORT_DIR, 'logo_' + safe)
    f.save(dest)
    s = _get_appdata()
    if s:
        s.rapport_logo_path = dest
        db.session.commit()
    return jsonify({'ok': True, 'path': dest})


# ---------------------------------------------------------------------------
# Report versions: GET /api/report/versions
# ---------------------------------------------------------------------------

@rapport_bp.route('/api/report/versions', methods=['GET'])
def rapport_versions():
    versions_dir = os.path.join(_RAPPORT_DIR, 'versions')
    if not os.path.exists(versions_dir):
        return jsonify({'versions': []})
    items = []
    for name in sorted(os.listdir(versions_dir), reverse=True):
        path = os.path.join(versions_dir, name)
        if os.path.isdir(path):
            files = os.listdir(path)
            items.append({'timestamp': name, 'files': files})
    return jsonify({'versions': items})


# ---------------------------------------------------------------------------
# Report validation: GET /api/report/validate
# ---------------------------------------------------------------------------

@rapport_bp.route('/api/report/validate', methods=['GET'])
def rapport_validate():
    s = _get_appdata()
    warnings = []
    errors = []

    if not s or not getattr(s, 'rapport_titel', ''):
        errors.append('Rapport titel is niet ingesteld')
    if not s or not getattr(s, 'rapport_auteur', ''):
        warnings.append('Auteur is niet ingesteld')

    findings = db_bevindingen.query.filter(db_bevindingen.status == 'final').all()
    if not findings:
        warnings.append('Geen findings met status "definitief"')

    templates = db_bevindingen_templates.query.all()
    tmpl_map = {str(t.id): t for t in templates}
    for f in findings:
        if not f.ref or str(f.ref) not in tmpl_map:
            warnings.append('Finding {} ({}) heeft geen geldige template'.format(f.id, f.naam or ''))
        if not f.basescore:
            tmpl = tmpl_map.get(str(f.ref))
            if not tmpl or not tmpl.basescore:
                warnings.append('Finding {} ({}) heeft geen CVSS score'.format(f.id, f.naam or ''))

    notes_count = db_notes.query.filter_by(rapport=True).count()
    if notes_count == 0:
        warnings.append('Geen notes geselecteerd voor het rapport')

    return jsonify({
        'ok': len(errors) == 0,
        'errors': errors,
        'warnings': warnings,
        'findings_count': len(findings),
    })


# ---------------------------------------------------------------------------
# Report metrics: GET /api/report/metrics
# ---------------------------------------------------------------------------

@rapport_bp.route('/api/report/metrics', methods=['GET'])
def rapport_metrics():
    findings = db_bevindingen.query.filter(db_bevindingen.status == 'final').all()
    templates = db_bevindingen_templates.query.all()
    tmpl_map = {str(t.id): t for t in templates}

    severity_dist = {'critical': 0, 'high': 0, 'medium': 0, 'low': 0, 'none': 0}
    remediation_dist = {'open': 0, 'in_progress': 0, 'fixed': 0, 'verified': 0, 'accepted': 0}
    owasp_dist = {}

    for f in findings:
        tmpl = tmpl_map.get(str(f.ref))
        bs = f.basescore or (tmpl.basescore if tmpl else '') or '0'
        sev, _ = _cvss_to_severity(bs)
        severity_dist[sev] = severity_dist.get(sev, 0) + 1

        rem = getattr(f, 'remediation_status', 'open') or 'open'
        remediation_dist[rem] = remediation_dist.get(rem, 0) + 1

        owasp_num = tmpl.owasp if tmpl else ''
        if owasp_num and str(owasp_num).isdigit():
            num = int(owasp_num)
            if 1 <= num <= 10:
                label = owasptop10[num - 1]
                owasp_dist[label] = owasp_dist.get(label, 0) + 1

    return jsonify({
        'total': len(findings),
        'severity': severity_dist,
        'remediation': remediation_dist,
        'owasp': owasp_dist,
    })


# ---------------------------------------------------------------------------
# Async report generation
# ---------------------------------------------------------------------------

def _async_generate(run_id, app, include_draft, include_checklists=''):
    """Worker die gen_rapport() uitvoert in een achtergrond-thread."""
    with _async_lock:
        _async_runs[run_id]['status'] = 'running'
    try:
        qs = 'include_draft=' + ('1' if include_draft else '0')
        if include_checklists:
            qs += '&include_checklists=' + include_checklists
        with app.test_request_context('/dashboard/report/generate?' + qs):
            gen_rapport()
        with _async_lock:
            _async_runs[run_id]['status'] = 'success'
            _async_runs[run_id]['finished_at'] = datetime.datetime.now().isoformat()
    except Exception as exc:
        with _async_lock:
            _async_runs[run_id]['status'] = 'failed'
            _async_runs[run_id]['finished_at'] = datetime.datetime.now().isoformat()
            _async_runs[run_id]['error'] = str(exc)


@rapport_bp.route('/api/report/generate', methods=['POST'])
def rapport_generate_async():
    """Start async rapport generatie. Retourneert run_id voor polling."""
    from flask import current_app
    run_id = str(_uuid.uuid4())
    include_draft = request.json.get('include_draft', False) if request.is_json else False
    include_checklists = request.json.get('include_checklists', '') if request.is_json else ''
    with _async_lock:
        _async_runs[run_id] = {
            'id': run_id,
            'status': 'queued',
            'started_at': datetime.datetime.now().isoformat(),
            'finished_at': None,
            'error': None,
        }
    t = threading.Thread(
        target=_async_generate,
        args=(run_id, current_app._get_current_object(), include_draft, include_checklists),
        daemon=True,
    )
    t.start()
    return jsonify({'ok': True, 'run_id': run_id})


@rapport_bp.route('/api/report/generate/<run_id>', methods=['GET'])
def rapport_generate_status(run_id):
    """Poll async generatie status."""
    with _async_lock:
        run = _async_runs.get(run_id)
    if not run:
        return jsonify({'error': 'run not found'}), 404
    return jsonify(run)
